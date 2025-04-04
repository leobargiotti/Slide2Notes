import requests
import json
import os
import time

from PyQt6.QtWidgets import QApplication
from pptx import Presentation
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.lib.units import inch
import fitz  # PyMuPDF
from docx import Document
import io
from PIL import Image
import base64

from sympy.physics.units import current


def send_request_to_api(prompt, max_retries=100):
    """
    Send a request to the Gemini API with a given prompt, retrying if the request fails due to a 429 error.

    Parameters:
    - prompt (str): The specific prompt to include in the request.
    - max_retries (int): Maximum number of retries for the request.

    Returns:
    - str: The response text or an error message.
    """
    url = os.getenv("GOOGLE_MODEL")
    api_key = os.getenv("API_KEY")

    model_config = {
        "temperature": 0,
    }

    headers = {
        "Content-Type": "application/json"
    }

    data = {
        "generation_config": model_config,
        "contents": [
            {
                "parts": [
                    {"text": prompt}
                ]
            }
        ]
    }

    retries = 0
    while retries <= max_retries:
        try:
            response = requests.post(f"{url}?key={api_key}", headers=headers, data=json.dumps(data))
            if response.status_code == 200:
                result = response.json()
                try:
                    return result['candidates'][0]['content']['parts'][0]['text'].replace("*", "")
                except (KeyError, IndexError):
                    raise Exception("Error: Unexpected response structure.")
            elif response.status_code == 429:
                retries += 1
                time.sleep(1)
            else:
                raise Exception(f"Error {response.status_code}: {response.text}")
        except Exception as e:
            raise Exception(e)
    raise Exception("Error: Maximum retries exceeded. Could not complete the request.")

def send_request_to_api_with_image(prompt, image_path, max_retries=1000):
    """
    Send a request to the Gemini API with a given prompt and image, retrying if the request fails due to a 429 error.

    Parameters:
    - prompt (str): The specific prompt to include in the request.
    - image_path (str): Path to the image file to analyze.
    - max_retries (int): Maximum number of retries for the request.

    Returns:
    - str: The response text or an error message.
    """
    # For Gemini Pro Vision, which supports image inputs
    url = os.getenv("GOOGLE_MODEL")
    api_key = os.getenv("API_KEY")


    model_config = {
        "temperature": 0,
    }

    headers = {
        "Content-Type": "application/json"
    }

    # Read and encode the image
    with open(image_path, "rb") as img_file:
        image_data = base64.b64encode(img_file.read()).decode("utf-8")

    # Structure the request for Gemini Pro Vision
    data = {
        "contents": [
            {
                "parts": [
                    {"text": prompt},
                    {
                        "inline_data": {
                            "mime_type": "image/png",
                            "data": image_data
                        }
                    }
                ]
            }
        ],
        "generation_config": model_config
    }

    retries = 0
    while retries <= max_retries:
        try:
            response = requests.post(f"{url}?key={api_key}", headers=headers, data=json.dumps(data))
            if response.status_code == 200:
                result = response.json()
                try:
                    return result['candidates'][0]['content']['parts'][0]['text'].replace("*", "")
                except (KeyError, IndexError):
                    raise Exception("Error: Unexpected response structure.")
            elif response.status_code == 429:
                retries += 1
                time.sleep(1)
            else:
                raise Exception(f"Error {response.status_code}: {response.text}")
        except Exception as e:
            raise Exception(str(e))
    raise Exception("Error: Maximum retries exceeded. Could not complete the request.")

def save_as_docx_file(output_path, summaries):
    """
    Saves a list of summaries as a DOCX file at the specified path.

    Parameters:
    - output_path (str): Path to the output file.
    - summaries (list): List of dictionaries, each containing a title and a content string.
    """
    doc = Document()
    for i, summary in enumerate(summaries):
        doc.add_heading(summary['title'], level=1)

        paragraphs = summary['content'].split('\n\n')
        for para in paragraphs:
            if para.strip():
                p = doc.add_paragraph(para)
                p.alignment = 3

        if i < len(summaries) - 1:
            doc.add_page_break()
    doc.save(output_path)

def save_as_pdf_file(output_path, summaries):

    """
    Save the summaries as a PDF file.

    Parameters:
    - output_path (str): The path where the PDF file will be saved.
    - summaries (list): A list of dictionaries containing the summaries. Each dictionary should have 'title' and 'content' keys.

    Example:
    summaries = [
        {'title': 'Section 1: File1.pdf', 'content': 'Summary content for File1.pdf'},
        {'title': 'Section 2: File2.pptx', 'content': 'Summary content for File2.pptx'}
    ]

    The PDF document will contain a section for each summary, with the title in a larger font and the content justified.
    """
    doc = SimpleDocTemplate(
        output_path,
        pagesize=A4,
        rightMargin=72,
        leftMargin=72,
        topMargin=72,
        bottomMargin=72
    )

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        'SectionTitle',
        parent=styles['Heading1'],
        fontSize=16,
        spaceAfter=12,
        textColor=colors.darkblue
    )
    normal_style = ParagraphStyle(
        'JustifiedText',
        parent=styles['Normal'],
        fontSize=11,
        leading=14,
        alignment=4
    )

    story = []

    for i, summary in enumerate(summaries):
        title = Paragraph(summary['title'], title_style)
        story.append(title)
        story.append(Spacer(1, 0.2 * inch))

        paragraphs = summary['content'].split('\n\n')
        for para in paragraphs:
            if para.strip():
                p = Paragraph(para.replace('\n', '\n'), normal_style)
                story.append(p)
                story.append(Spacer(1, 0.1 * inch))

        if i < len(summaries) - 1:
            story.append(PageBreak())

    doc.build(story)


def extract_text_from_pptx(file_path, progress, current_page_progress):
    """
    Extract text from a PowerPoint (.pptx) file, including slide content and notes.

    Parameters:
    - file_path (str): The path to the .pptx file from which text will be extracted.
    - progress (QProgressBar): A progress bar to update when extracting text.
    - current_page_progress (int): The current page number.

    Returns:
    - str: A string containing the extracted text from the slides and their notes,
           with slide content separated by newlines and notes prefixed with "Note:".
    - int: The updated current page number.
    """
    text = ""
    presentation = Presentation(file_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
        # Extract slide notes
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            for shape in notes_slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += f"\nNote: {shape.text}"
        current_page_progress += 1
        progress.setValue(current_page_progress)
        QApplication.processEvents()
    return text, current_page_progress

def extract_text_and_images_from_pptx(file_path, progress, current_page_progress):
    """
    Extract text and images from a PowerPoint (.pptx) file, including slide content, notes, and AI-generated descriptions of images.

    Parameters:
    - file_path (str): The path to the .pptx file from which text and images will be extracted.
    - progress (QProgressBar): A progress bar to update when extracting text and images.
    - current_page_progress (int): The current page number.

    Returns:
    - str: A string containing the extracted text from the slides and their notes, with slide content separated by newlines and notes prefixed with Note:, and AI-generated image descriptions prefixed with Image Description:.
    - int: The updated current page number.
    """
    text = ""
    presentation = Presentation(file_path)

    for i, slide in enumerate(presentation.slides):
        text += f"\n\n--- Slide {i + 1} ---\n"

        # Extract text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text += shape.text.strip() + "\n"

            # Process images
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    # Extract image
                    image = shape.image
                    image_bytes = image.blob

                    # Save image to memory buffer
                    image_buffer = io.BytesIO(image_bytes)
                    pil_image = Image.open(image_buffer)

                    if pil_image.mode == 'CMYK':
                        pil_image = pil_image.convert('RGB')

                    # Generate a temporary file path to save the image
                    temp_image_path = f"temp_image_slide_{i + 1}_{shape.name}.png"
                    pil_image.save(temp_image_path)

                    # Generate description using AI
                    try:
                        # Modified send_request_to_api that handles images
                        image_description = send_request_to_api_with_image(
                            prompt="Describe this image in 2-3 sentences. Focus on the main elements visible in the image.",
                            image_path=temp_image_path,  # Or use base64_image if API accepts base64
                        )
                        text += f"\n[Image Description: {image_description}]\n"

                    except Exception as e:
                        print(f"Error generating image {temp_image_path} in {os.path.basename(file_path)} with description: {str(e)}")

                    # Cleanup temporary file
                    if os.path.exists(temp_image_path):
                        os.remove(temp_image_path)

                except Exception as e:
                    print(f"Error generating image at slide {i + 1} in {os.path.basename(file_path)} with description: {str(e)}")

        # Extract slide notes
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            for shape in notes_slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += f"\nNote: {shape.text}"

        current_page_progress += 1
        progress.setValue(current_page_progress)
        QApplication.processEvents()

    return text, current_page_progress

def extract_text_from_pdf(file_path, progress, current_page_progress):
    """
    Extract text from a PDF file, including page content and annotations.

    Parameters:
    - file_path (str): The path to the PDF file from which text will be extracted.
    - progress (QProgressBar): A progress bar to update when extracting text.
    - current_page_progress (int): The current page number.

    Returns:
    - str: A string containing the extracted text from the PDF file,
           with annotations prefixed with "Note:".
    - int: The updated current page number.
    """
    text = ""
    pdf_document = fitz.open(file_path)
    for page in pdf_document:
        text += page.get_text()
        # Extract annotations
        for annot in page.annots():
            text += f"\nNote: {annot.info['content']}"
        current_page_progress += 1
        progress.setValue(current_page_progress)
        QApplication.processEvents()
    pdf_document.close()
    return text, current_page_progress

def extract_text_and_images_from_pdf(file_path, progress, current_page_progress):
    """
    Extract text and images from a PDF file, including page content, annotations, and images.

    Parameters:
    - file_path (str): The path to the PDF file from which text and images will be extracted.
    - progress (QProgressBar): A progress bar to update when extracting text and images.
    - current_page_progress (int): The current page number.

    Returns:
    - str: A string containing the extracted text from the PDF file, with annotations prefixed with Note: and AI-generated image descriptions prefixed with Image Description:.
    - int: The updated current page number.
    """
    text = ""
    pdf_document = fitz.open(file_path)

    for page_num, page in enumerate(pdf_document):
        text += f"\n\n--- Page {page_num + 1} ---\n"

        # Extract text from page
        page_text = page.get_text()
        if page_text.strip():
            text += page_text.strip() + "\n"

        # Extract images
        image_list = page.get_images(full=True)

        for img_index, img_info in enumerate(image_list):
            try:
                # Get the image
                xref = img_info[0]
                base_image = pdf_document.extract_image(xref)
                image_bytes = base_image["image"]

                # Create a PIL image
                image_buffer = io.BytesIO(image_bytes)
                pil_image = Image.open(image_buffer)

                # Handle CMYK images
                if pil_image.mode == 'CMYK':
                    pil_image = pil_image.convert('RGB')

                # Save image to temporary file
                temp_image_path = f"temp_image_page_{page_num + 1}_{img_index}.png"
                pil_image.save(temp_image_path)

                # Generate description using AI
                try:
                    # Using the new function that can handle images
                    image_description = send_request_to_api_with_image(
                        prompt="Describe this image in 2-3 sentences. Focus on the main elements visible in the image.",
                        image_path=temp_image_path
                    )
                    text += f"\n[Image Description: {image_description}]\n"
                except Exception as e:
                    print(f"Error generating image on page {page_num + 1} in {os.path.basename(file_path)} with description: {str(e)}")

                # Cleanup temporary file
                if os.path.exists(temp_image_path):
                    os.remove(temp_image_path)

            except Exception as e:
                print(f"Error processing image on page {page_num + 1} in {os.path.basename(file_path)}: {str(e)}")

        # Extract annotations
        for annot in page.annots():
            if "content" in annot.info and annot.info["content"].strip():
                text += f"\nNote: {annot.info['content']}\n"

        current_page_progress += 1
        progress.setValue(current_page_progress)
        QApplication.processEvents()

    pdf_document.close()
    return text, current_page_progress


def create_summary_prompt(text, target_language):
    """
    Create a fully expanded and cohesive textual version of the given content in the specified target language.
    The output must preserve all information while transforming it into a natural, readable narrative.

    Parameters:
    - text (str): The full content, including main text, notes, and image descriptions.
    - target_language (str): The language in which the expanded text should be provided.

    Returns:
    - str: A formatted prompt string for generating the expanded text.
    """
    return f"""Please rewrite the following content in {target_language} as a fully expanded, cohesive, and detailed narrative.

                CRITICAL INSTRUCTION: You MUST include EVERYTHING from:
                1. The main text
                2. All notes (marked with "Note:")
                3. All image descriptions (marked with "Image Description:")
                4. All formulas (mathematical or otherwise), exactly as provided.

                **DO NOT include any structural references like "Page 1" or "Image 1". These references should not appear in the final text.**

                Your output MUST fully retain **every single detail** provided. Do NOT summarize or omit any information, no matter how minor it seems. The final narrative should include **all information** from the text, notes, image descriptions, and formulas exactly as provided, **with no details left out**.

                When describing visual elements such as graphs, charts, or diagrams, ensure that **every detail** of the description is fully explained, and their significance or meaning is clearly conveyed. **Do not abbreviate or condense** the descriptions of images—fully explain what they represent and their relevance to the overall content.

                When formulas are included, you must **transcribe them exactly as they appear**, and explain their significance and how they fit into the broader context of the content. **Do not leave any formulas out** or reduce their complexity.

                Your narrative should flow seamlessly, as if all the information was originally part of a cohesive document. The content should be integrated smoothly into one continuous narrative without separating the different components.

                Text to expand:
                {text}

                The result should be detailed, thorough, and well-structured, resembling an informative article or lecture that seamlessly incorporates every detail from all sources without leaving anything out or overly condensing any part. Avoid bullet points and ensure the final text is rich in information and clarity."""